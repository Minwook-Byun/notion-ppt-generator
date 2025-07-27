#!/usr/bin/env python3
"""
Enhanced PowerPoint MCP Server with Template Clone & Update + Notion Integration
Core functionality: Template cloning, smart content update, and Notion data integration
"""
import os
import json
import datetime
import shutil
import re
from pathlib import Path
from typing import Optional, Dict, List, Any
from mcp.server.fastmcp import FastMCP
from pptx import Presentation
from pptx.util import Pt, Inches
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
# from pptx.enum.table import MSO_TABLE_ALIGNMENT  # Not available in this version
import tempfile
from copy import deepcopy

# Create FastMCP server
mcp = FastMCP("Enhanced PowerPoint MCP Server with Template Clone & Update + Notion Integration")

# Global presentation object
current_presentation = None
current_filename = None

# ═══════════════════════════════════════════════════════════════════
# TEMPLATE AND DIRECTORY CONFIGURATION
# ═══════════════════════════════════════════════════════════════════

# Common template paths (accessible to all users)
TEMPLATE_PATHS = {
    'common_templates': Path("C:/Templates/PowerPoint"),
    'public_documents': Path.home() / "Documents" / "PowerPoint Templates",
    'desktop_templates': Path.home() / "Desktop" / "Templates",
    'shared_drive': Path("//shared/templates/powerpoint"),
}

# Setup save directories
PRESENTATIONS_DIR = Path.home() / "Desktop" / "MyPPT"
PRESENTATIONS_DIR.mkdir(exist_ok=True)

# Setup temp directory
TEMP_DIR = Path(tempfile.gettempdir()) / "mcp_powerpoint"
TEMP_DIR.mkdir(exist_ok=True)

# Template registry
template_registry = {}

def discover_templates():
    """Discover available PowerPoint templates from common locations"""
    global template_registry
    template_registry = {}
    
    for location_name, path in TEMPLATE_PATHS.items():
        if path.exists():
            templates = list(path.glob("*.pptx")) + list(path.glob("*.potx"))
            for template_path in templates:
                template_name = template_path.stem
                template_registry[template_name] = {
                    'path': str(template_path),
                    'location': location_name,
                    'name': template_name,
                    'extension': template_path.suffix
                }
    
    return template_registry

def load_template_presentation(template_name: str) -> Optional[Presentation]:
    """Load a PowerPoint template by name"""
    if template_name not in template_registry:
        return None
    
    template_info = template_registry[template_name]
    template_path = template_info['path']
    
    try:
        return Presentation(template_path)
    except Exception as e:
        print(f"Error loading template {template_name}: {e}")
        return None

def update_presentation_with_smart_text(presentation: Presentation, chapter_text: str = "", 
                                      title_text: str = "", contents_text: str = "") -> Dict:
    """
    Smart text update for presentations.
    Finds placeholder text in template and replaces with actual content
    """
    modified_count = 0
    
    try:
        for slide_idx, slide in enumerate(presentation.slides):
            for shape in slide.shapes:
                if not shape.has_text_frame:
                    continue
                text = shape.text_frame.text.strip().lower()
                
                # Chapter text update
                if any(keyword in text for keyword in ['chapter', 'chap']):
                    if chapter_text:
                        shape.text_frame.text = chapter_text
                        # Font settings
                        for paragraph in shape.text_frame.paragraphs:
                            paragraph.font.name = "Pretendard"
                            paragraph.font.size = Pt(18)
                        modified_count += 1
                
                # Title text update
                elif any(keyword in text for keyword in ['title']):
                    if title_text:
                        shape.text_frame.text = title_text
                        # Font settings
                        for paragraph in shape.text_frame.paragraphs:
                            paragraph.font.name = "Pretendard"
                            paragraph.font.size = Pt(24)
                        modified_count += 1
                        
                # Contents text update
                elif any(keyword in text for keyword in ['contents', 'content']):
                    if contents_text:
                        shape.text_frame.text = contents_text
                        # Font settings
                        for paragraph in shape.text_frame.paragraphs:
                            paragraph.font.name = "Pretendard"
                            paragraph.font.size = Pt(14)
                        modified_count += 1
        
        return {
            "status": "success",
            "modified_count": modified_count,
            "error_message": None
        }
    except Exception as e:
        return {
            "status": "failure",
            "modified_count": 0,
            "error_message": str(e)
        }

# ═══════════════════════════════════════════════════════════════════
# NOTION INTEGRATION HELPER FUNCTIONS
# ═══════════════════════════════════════════════════════════════════

def fetch_notion_page(notion_url: str) -> Dict[str, Any]:
    """
    Notion 페이지 전체 내용 가져오기
    실제로는 Notion:fetch 커넥터 함수 호출
    """
    # TODO: 실제 구현에서는 Notion 커넥터 사용
    # return Notion.fetch(notion_url)
    return {}

def search_notion_database(query: str, database_url: str = None) -> List[Dict[str, Any]]:
    """
    Notion 데이터베이스 검색
    실제로는 Notion:search 커넥터 함수 호출
    """
    # TODO: 실제 구현에서는 Notion 커넥터 사용
    # return Notion.search(query, database_url)
    return []

def extract_basic_settings(notion_content: str) -> Dict[str, Any]:
    """
    Notion 페이지에서 기본 설정 추출
    
    기대 구조:
    ## 🔧 기본 설정
    **프로젝트명:** K-Camp 제주 3기 소개자료
    **템플릿:** MYSC_Sample_Template
    **폰트:** Pretendard
    **총 슬라이드 수:** 5
    """
    settings = {}
    
    # 기본 설정 섹션 찾기
    basic_section = re.search(r'## 🔧 기본 설정(.*?)(?=##|$)', notion_content, re.DOTALL)
    if not basic_section:
        return {}
    
    content = basic_section.group(1)
    
    # 각 설정 항목 파싱
    patterns = {
        'project_name': r'\*\*프로젝트명:\*\* (.+)',
        'template': r'\*\*템플릿:\*\* (.+)',
        'font': r'\*\*폰트\*\*?\s*[:：]\s*(.+)',
        'diagram_type': r'\*\*다이어그램 타입:\*\* (.+)',
        'total_slides': r'\*\*총 슬라이드 수:\*\* (\d+)'
    }
    
    for key, pattern in patterns.items():
        match = re.search(pattern, content)
        if match:
            value = match.group(1).strip()
            if key == 'total_slides':
                settings[key] = int(value)
            else:
                settings[key] = value
    
    return settings

def extract_slide_configurations(notion_url: str) -> List[Dict[str, Any]]:
    """
    슬라이드 구성표 데이터베이스에서 설정 추출
    
    기대 데이터베이스 구조:
    | 슬라이드 | Chapter | Title | Contents | Layout_Type | 특별요구사항 |
    """
    try:
        # 실제로는 슬라이드 구성 데이터베이스 검색
        page_content = fetch_notion_page(notion_url)
        if not page_content:
            return []
        
        # 슬라이드 구성표 데이터베이스 URL 추출
        db_pattern = r'슬라이드 구성.*?https://www\.notion\.so/([a-f0-9]+)'
        db_match = re.search(db_pattern, page_content.get('text', ''))
        
        if db_match:
            db_id = db_match.group(1)
            slide_configs = search_notion_database("", f"collection://{db_id}")
            return slide_configs
        
        return []
        
    except Exception as e:
        print(f"슬라이드 구성 추출 오류: {e}")
        return []

def extract_table_data(notion_url: str) -> List[Dict[str, Any]]:
    """
    표 데이터 데이터베이스에서 표 정보 추출
    
    기대 데이터베이스 구조:
    | Table_ID | Parent_Slide | Row | Column | Cell_Value | Header_Type |
    """
    try:
        page_content = fetch_notion_page(notion_url)
        if not page_content:
            return []
        
        # 표 데이터 데이터베이스 URL 추출
        db_pattern = r'표 데이터.*?https://www\.notion\.so/([a-f0-9]+)'
        db_match = re.search(db_pattern, page_content.get('text', ''))
        
        if db_match:
            db_id = db_match.group(1)
            table_data = search_notion_database("", f"collection://{db_id}")
            return table_data
        
        return []
        
    except Exception as e:
        print(f"표 데이터 추출 오류: {e}")
        return []

def extract_style_guide(notion_url: str) -> Dict[str, Any]:
    """Notion 페이지에서 스타일 가이드 추출"""
    try:
        page_content = fetch_notion_page(notion_url)
        if not page_content:
            return get_default_style_guide()
        
        content_text = page_content.get('text', '')
        
        # 스타일 가이드 섹션 찾기
        style_section = re.search(r'## 🎨 스타일 가이드(.*?)(?=##|$)', content_text, re.DOTALL)
        if not style_section:
            return get_default_style_guide()
        
        return parse_style_guide_content(style_section.group(1))
        
    except Exception as e:
        print(f"스타일 가이드 추출 오류: {e}")
        return get_default_style_guide()

def get_default_style_guide() -> Dict[str, Any]:
    """기본 스타일 가이드 반환"""
    return {
        'colors': {
            'main': '#1E3A8A',
            'accent': '#F97316', 
            'background': '#F8FAFC'
        },
        'fonts': {
            'title': {'size': 24, 'bold': True},
            'body': {'size': 14, 'bold': False},
            'caption': {'size': 12, 'bold': False}
        },
        'layout': {
            'margin': Inches(2/2.54),
            'spacing': Inches(1/2.54)
        }
    }

def parse_style_guide_content(content: str) -> Dict[str, Any]:
    """스타일 가이드 콘텐츠 파싱"""
    style_guide = get_default_style_guide()
    
    # 색상 팔레트 파싱
    color_patterns = {
        'main': r'\*\*메인 컬러:\*\* (#[A-Fa-f0-9]{6})',
        'accent': r'\*\*강조 컬러:\*\* (#[A-Fa-f0-9]{6})',
        'background': r'\*\*배경 컬러:\*\* (#[A-Fa-f0-9]{6})'
    }
    
    for key, pattern in color_patterns.items():
        match = re.search(pattern, content)
        if match:
            style_guide['colors'][key] = match.group(1)
    
    # 폰트 설정 파싱
    font_patterns = {
        'title': r'\*\*제목:\*\* [^,]*,?\s*(\d+)pt',
        'body': r'\*\*본문:\*\* [^,]*,?\s*(\d+)pt',
        'caption': r'\*\*캡션:\*\* [^,]*,?\s*(\d+)pt'
    }
    
    for key, pattern in font_patterns.items():
        match = re.search(pattern, content)
        if match:
            style_guide['fonts'][key]['size'] = int(match.group(1))
            # 굵게 여부 확인
            bold_check = re.search(rf'\*\*{key}:\*\* 굵게', content)
            style_guide['fonts'][key]['bold'] = bool(bold_check)
    
    return style_guide

def parse_notion_color(color_str: str) -> RGBColor:
    """Notion 색상 문자열을 RGBColor로 변환"""
    if color_str.startswith('#'):
        color_str = color_str[1:]
    
    r = int(color_str[0:2], 16)
    g = int(color_str[2:4], 16) 
    b = int(color_str[4:6], 16)
    
    return RGBColor(r, g, b)

def organize_table_data(table_entries: List[Dict]) -> Dict[str, Dict]:
    """표 데이터를 슬라이드별로 정리"""
    tables_by_slide = {}
    
    for entry in table_entries:
        parent_slide = entry.get('Parent_Slide', '1')
        table_id = entry.get('Table_ID', 'DEFAULT')
        
        key = f"slide_{parent_slide}_{table_id}"
        
        if key not in tables_by_slide:
            tables_by_slide[key] = {
                'slide_number': parent_slide,
                'table_id': table_id,
                'cells': [],
                'max_row': 0,
                'max_col': 0
            }
        
        row = int(entry.get('Row', 1))
        col = int(entry.get('Column', 1))
        
        tables_by_slide[key]['max_row'] = max(tables_by_slide[key]['max_row'], row)
        tables_by_slide[key]['max_col'] = max(tables_by_slide[key]['max_col'], col)
        
        tables_by_slide[key]['cells'].append({
            'row': row,
            'col': col,
            'value': entry.get('Cell_Value', ''),
            'is_header': entry.get('Header_Type') == 'column_header'
        })
    
    return tables_by_slide

# ═══════════════════════════════════════════════════════════════════
# TEMPLATE DISCOVERY AND MANAGEMENT TOOLS
# ═══════════════════════════════════════════════════════════════════

@mcp.tool()
def scan_templates() -> str:
    """Scan and discover available PowerPoint templates"""
    try:
        templates = discover_templates()
        
        if not templates:
            return """No templates found.
            
Searched locations:
""" + "\n".join([f"   • {name}: {path}" for name, path in TEMPLATE_PATHS.items()])
        
        result = [f"Found {len(templates)} templates:\n"]
        
        # Group by location
        by_location = {}
        for template_name, info in templates.items():
            location = info['location']
            if location not in by_location:
                by_location[location] = []
            by_location[location].append(info)
        
        for location, template_list in by_location.items():
            result.append(f"{location.replace('_', ' ').title()}:")
            for template in template_list:
                result.append(f"   {template['name']}{template['extension']}")
            result.append("")
        
        return "\n".join(result)
        
    except Exception as e:
        return f"Error scanning templates: {str(e)}"

@mcp.tool()
def list_available_templates() -> str:
    """List all currently available templates"""
    if not template_registry:
        discover_templates()
    
    if not template_registry:
        return "No templates available. Run 'scan_templates' first."
    
    result = [f"Available templates ({len(template_registry)}):\n"]
    
    for template_name, info in sorted(template_registry.items()):
        result.append(f"{template_name}")
        result.append(f"   Location: {info['location']}")
        result.append(f"   Path: {info['path']}")
        result.append("")
    
    return "\n".join(result)

# ═══════════════════════════════════════════════════════════════════
# SLIDE DUPLICATION TOOLS
# ═══════════════════════════════════════════════════════════════════

@mcp.tool()
def duplicate_slide(slide_number: int = 1, new_title: str = "", new_content: str = "", new_chapter: str = "") -> str:
    """
    지정된 슬라이드를 복제하여 새로운 슬라이드 생성
    
    Args:
        slide_number: 복제할 슬라이드 번호 (기본값: 1)
        new_title: 새 슬라이드의 제목
        new_content: 새 슬라이드의 내용  
        new_chapter: 새 슬라이드의 챕터
    
    Returns:
        복제 결과 메시지
    """
    global current_presentation
    
    if not current_presentation:
        return "No presentation is currently open"
    
    try:
        # 복제할 슬라이드 선택 (인덱스는 0부터 시작)
        if slide_number < 1 or slide_number > len(current_presentation.slides):
            return f"Invalid slide number. Must be between 1 and {len(current_presentation.slides)}"
        
        source_slide = current_presentation.slides[slide_number - 1]
        
        # 슬라이드 레이아웃 가져오기
        slide_layout = source_slide.slide_layout
        
        # 새 슬라이드 추가
        new_slide = current_presentation.slides.add_slide(slide_layout)
        
        # 소스 슬라이드의 모든 shape 복제
        for shape in source_slide.shapes:
            try:
                # shape 요소를 깊은 복사
                new_shape_element = deepcopy(shape.element)
                new_slide.shapes._spTree.insert_element_before(
                    new_shape_element, 'p:extLst'
                )
            except Exception as shape_error:
                print(f"Warning: Could not copy shape: {shape_error}")
                continue
        
        # 새로운 내용으로 업데이트
        slide_count = len(current_presentation.slides)
        updated_elements = 0
        
        if new_title or new_content or new_chapter:
            # 텍스트 업데이트
            for shape in new_slide.shapes:
                if not shape.has_text_frame:
                    continue
                    
                text = shape.text.lower()
                
                # 챕터 업데이트
                if new_chapter and ('chapter' in text or 'chap' in text):
                    shape.text = new_chapter
                    for paragraph in shape.text_frame.paragraphs:
                        paragraph.font.name = "Pretendard"
                        paragraph.font.size = Pt(18)
                    updated_elements += 1
                    
                # 제목 업데이트  
                elif new_title and ('title' in text or len(text) < 50):
                    shape.text = new_title
                    for paragraph in shape.text_frame.paragraphs:
                        paragraph.font.name = "Pretendard"
                        paragraph.font.size = Pt(24)
                    updated_elements += 1
                    
                # 내용 업데이트
                elif new_content and ('content' in text or len(text) > 50):
                    shape.text = new_content
                    for paragraph in shape.text_frame.paragraphs:
                        paragraph.font.name = "Pretendard"
                        paragraph.font.size = Pt(14)
                    updated_elements += 1
        
        return f"Successfully duplicated slide {slide_number} -> slide {slide_count}" + \
               (f" with {updated_elements} text elements updated" if new_title or new_content or new_chapter else "")
        
    except Exception as e:
        return f"Error duplicating slide: {str(e)}"

@mcp.tool()
def add_slide_copy_first(content: str = "", title: str = "", chapter: str = "", layout_index: int = 1) -> str:
    """
    첫 번째 슬라이드를 복제하여 새로운 슬라이드를 추가 (기존 add_slide 대체용)
    
    Args:
        content: 슬라이드 내용
        title: 슬라이드 제목
        chapter: 슬라이드 챕터
        layout_index: 레이아웃 인덱스 (사용되지 않음, 호환성용)
    
    Returns:
        슬라이드 추가 결과
    """
    return duplicate_slide(1, title, content, chapter)

@mcp.tool()  
def duplicate_and_update_slides(count: int = 4, slide_data: List[Dict] = None) -> str:
    """
    첫 번째 슬라이드를 여러 번 복제하고 각각 다른 내용으로 업데이트
    
    Args:
        count: 복제할 슬라이드 개수
        slide_data: 각 슬라이드의 데이터 [{"title": "", "content": "", "chapter": ""}, ...]
    
    Returns:
        복제 및 업데이트 결과
    """
    global current_presentation
    
    if not current_presentation:
        return "No presentation is currently open"
    
    if not slide_data:
        slide_data = []
    
    results = []
    
    for i in range(count):
        data = slide_data[i] if i < len(slide_data) else {}
        title = data.get('title', f'슬라이드 {i+2}')
        content = data.get('content', f'내용 {i+2}')
        chapter = data.get('chapter', f'Chapter {i+2}')
        
        result = duplicate_slide(1, title, content, chapter)
        results.append(f"Slide {i+2}: {result}")
    
    return "\n".join(results)

# ═══════════════════════════════════════════════════════════════════
# NOTION INTEGRATION TOOLS (NEW)
# ═══════════════════════════════════════════════════════════════════

@mcp.tool()
def insert_table_from_data(slide_number: int, table_data: Dict, style: Dict) -> str:
    """Notion 표 데이터를 PPT 표로 삽입"""
    global current_presentation
    
    if current_presentation is None:
        return "❌ 열린 프레젠테이션이 없습니다."
    
    try:
        if slide_number > len(current_presentation.slides):
            return f"❌ 슬라이드 {slide_number}이 존재하지 않습니다."
        
        slide = current_presentation.slides[slide_number - 1]
        
        max_row = table_data['max_row']
        max_col = table_data['max_col']
        
        # 테이블 위치 및 크기 설정
        left = style.get('layout', {}).get('margin', Inches(1))
        top = Inches(3.5)  # 제목과 contents 아래
        width = Inches(8)
        height = Inches(0.5) * max_row
        
        # 테이블 추가
        table_shape = slide.shapes.add_table(max_row, max_col, left, top, width, height)
        table = table_shape.table
        
        # 2차원 배열로 데이터 정리
        cell_matrix = [['' for _ in range(max_col)] for _ in range(max_row)]
        header_matrix = [[False for _ in range(max_col)] for _ in range(max_row)]
        
        for cell_data in table_data['cells']:
            row_idx = cell_data['row'] - 1
            col_idx = cell_data['col'] - 1
            if 0 <= row_idx < max_row and 0 <= col_idx < max_col:
                cell_matrix[row_idx][col_idx] = cell_data['value']
                header_matrix[row_idx][col_idx] = cell_data['is_header']
        
        # 테이블 셀에 데이터 입력 및 스타일 적용
        for row_idx in range(max_row):
            for col_idx in range(max_col):
                cell = table.cell(row_idx, col_idx)
                cell.text = cell_matrix[row_idx][col_idx]
                
                if header_matrix[row_idx][col_idx]:
                    # 헤더 스타일
                    fill = cell.fill
                    fill.solid()
                    fill.fore_color.rgb = parse_notion_color(style['colors']['main'])
                    
                    for paragraph in cell.text_frame.paragraphs:
                        for run in paragraph.runs:
                            run.font.color.rgb = RGBColor(255, 255, 255)
                            run.font.bold = True
                            run.font.size = Pt(style['fonts']['caption']['size'])
                else:
                    # 일반 셀 스타일
                    for paragraph in cell.text_frame.paragraphs:
                        for run in paragraph.runs:
                            run.font.size = Pt(style['fonts']['body']['size'])
                            run.font.color.rgb = RGBColor(0, 0, 0)
        
        # 테이블 중앙 정렬
        slide_width = slide.shapes[0].width if slide.shapes else Inches(10)
        table_shape.left = int((slide_width - table_shape.width) / 2)
        
        return f"✅ 슬라이드 {slide_number}에 {max_row}x{max_col} 테이블 추가 완료"
        
    except Exception as e:
        return f"❌ 테이블 삽입 오류: {str(e)}"

@mcp.tool()
def apply_color_theme(color_palette: Dict[str, str]) -> str:
    """Notion 색상 가이드를 PPT에 적용"""
    global current_presentation
    
    if current_presentation is None:
        return "❌ 열린 프레젠테이션이 없습니다."
    
    try:
        main_color = parse_notion_color(color_palette.get('main', '#1E3A8A'))
        accent_color = parse_notion_color(color_palette.get('accent', '#F97316'))
        
        theme_applied_count = 0
        
        for slide in current_presentation.slides:
            # 제목 색상 변경
            if hasattr(slide.shapes, 'title') and slide.shapes.title:
                title_shape = slide.shapes.title
                if title_shape.has_text_frame:
                    for paragraph in title_shape.text_frame.paragraphs:
                        for run in paragraph.runs:
                            run.font.color.rgb = main_color
                            run.font.bold = True
                    theme_applied_count += 1
            
            # Chapter 텍스트 강조 색상 적용
            for shape in slide.shapes:
                if hasattr(shape, 'text_frame') and shape.has_text_frame:
                    text = shape.text_frame.text
                    if 'Chapter' in text:
                        for paragraph in shape.text_frame.paragraphs:
                            for run in paragraph.runs:
                                if 'Chapter' in run.text:
                                    run.font.color.rgb = accent_color
                                    run.font.bold = True
        
        return f"✅ 색상 테마 적용 완료 ({theme_applied_count}개 슬라이드, 메인: {color_palette['main']}, 강조: {color_palette['accent']})"
        
    except Exception as e:
        return f"❌ 색상 테마 적용 오류: {str(e)}"

@mcp.tool()
def auto_generate_from_notion_url(notion_url: str) -> str:
    """
    Notion URL 하나로 완전 자동 PPT 생성
    실시간으로 Notion 데이터를 가져와서 처리
    """
    try:
        results = []
        
        # 1. Notion 페이지에서 기본 설정 추출
        page_content = fetch_notion_page(notion_url)
        if not page_content:
            return "❌ Notion 페이지를 가져올 수 없습니다. Notion 커넥터가 활성화되어 있는지 확인하세요."
        
        basic_settings = extract_basic_settings(page_content.get('text', ''))
        if not basic_settings:
            return "❌ 기본 설정을 찾을 수 없습니다. Notion 페이지 구조를 확인하세요."
        
        results.append(f"📋 기본 설정 추출 완료: {basic_settings.get('project_name', 'Unknown')}")
        
        # 2. 스타일 가이드 추출
        style_guide = extract_style_guide(notion_url)
        results.append(f"🎨 스타일 가이드 추출 완료")
        
        # 3. 프레젠테이션 생성
        template_name = basic_settings.get('template', 'MYSC_Sample_Template')
        project_name = basic_settings.get('project_name', 'New Presentation')
        
        create_result = create_presentation_from_template(template_name, project_name)
        if "오류" in create_result or "not found" in create_result:
            return f"❌ 템플릿 생성 실패: {create_result}"
        results.append(create_result)
        
        # 4. 슬라이드 구성 데이터 추출 및 생성
        slide_configs = extract_slide_configurations(notion_url)
        if slide_configs:
            for i, config in enumerate(slide_configs):
                if i == 0:
                    # 첫 번째 슬라이드 업데이트
                    result = update_specific_slide_text(
                        slide_number=1,
                        chapter=config.get('Chapter', ''),
                        title=config.get('Title', ''), 
                        contents=config.get('Contents', '')
                    )
                else:
                    # 나머지 슬라이드 복제 생성
                    result = duplicate_slide(
                        slide_number=1,
                        new_chapter=config.get('Chapter', ''),
                        new_title=config.get('Title', ''),
                        new_content=config.get('Contents', '')
                    )
                results.append(result)
        else:
            results.append("⚠️ 슬라이드 구성 데이터를 찾을 수 없습니다.")
        
        # 5. 표 데이터 처리
        table_data_raw = extract_table_data(notion_url)
        if table_data_raw:
            tables_organized = organize_table_data(table_data_raw)
            
            for table_key, table_info in tables_organized.items():
                slide_num = int(table_info['slide_number'])
                table_result = insert_table_from_data(slide_num, table_info, style_guide)
                results.append(table_result)
        else:
            results.append("ℹ️ 표 데이터가 없습니다.")
        
        # 6. 색상 테마 적용
        color_result = apply_color_theme(style_guide['colors'])
        results.append(color_result)
        
        # 7. 저장
        filename = f"{project_name}_자동생성"
        save_result = save_presentation(filename)
        results.append(save_result)
        
        # 8. 결과 종합
        success_msg = f"""
🎉 Notion 기반 PPT 자동 생성 완료!

📊 처리 결과:
{chr(10).join(f"  {result}" for result in results)}

📁 URL: {notion_url}
💾 저장: {filename}.pptx
        """
        
        return success_msg.strip()
        
    except Exception as e:
        return f"❌ 자동 생성 오류: {str(e)}"

# ═══════════════════════════════════════════════════════════════════
# CORE FEATURE: TEMPLATE CLONE AND CONTENT UPDATE
# ═══════════════════════════════════════════════════════════════════

@mcp.tool()
def clone_template_and_update(template_name: str, chapter: str = "", title: str = "", 
                             contents: str = "", new_filename: str = "") -> str:
    """
    CORE FEATURE: Clone template and update chapter, title, contents
    """
    global current_presentation, current_filename
    
    if not template_registry:
        discover_templates()
    
    if template_name not in template_registry:
        available = ", ".join(template_registry.keys()) if template_registry else "None"
        return f"Template '{template_name}' not found.\nAvailable templates: {available}"
    
    try:
        # 1. Load template
        current_presentation = load_template_presentation(template_name)
        
        if current_presentation is None:
            return f"Failed to load template '{template_name}'"
        
        # 2. Smart text update
        update_result = update_presentation_with_smart_text(
            current_presentation, chapter, title, contents
        )
        
        if update_result["status"] != "success":
            return f"Failed to update content: {update_result['error_message']}"
        
        # 3. Set filename
        if new_filename:
            current_filename = new_filename if new_filename.endswith('.pptx') else f"{new_filename}.pptx"
        else:
            timestamp = datetime.datetime.now().strftime("%Y%m%d_%H%M%S")
            current_filename = f"{template_name}_updated_{timestamp}.pptx"
        
        template_info = template_registry[template_name]
        
        return f"""Template cloned and updated successfully!

Template: {template_name}
Source: {template_info['location']}
Modified elements: {update_result['modified_count']}
Total slides: {len(current_presentation.slides)}

Updated content:
   Chapter: {chapter if chapter else '(unchanged)'}
   Title: {title if title else '(unchanged)'}
   Contents: {contents[:50] + '...' if len(contents) > 50 else contents if contents else '(unchanged)'}

Ready to save as: {current_filename}
Use 'save_presentation()' to save the updated presentation!"""
        
    except Exception as e:
        return f"Error cloning and updating template: {str(e)}"

@mcp.tool()
def update_specific_slide_text(slide_number: int, chapter: str = "", title: str = "", 
                              contents: str = "") -> str:
    """
    Update text in a specific slide
    """
    global current_presentation
    
    if current_presentation is None:
        return "No presentation open. Please create or load a presentation first."
    
    if slide_number < 1 or slide_number > len(current_presentation.slides):
        return f"Invalid slide number. Presentation has {len(current_presentation.slides)} slides."
    
    try:
        slide = current_presentation.slides[slide_number - 1]
        modified_count = 0
        
        for shape in slide.shapes:
            if not shape.has_text_frame:
                continue
            text = shape.text_frame.text.strip().lower()
            
            # Chapter text update
            if any(keyword in text for keyword in ['chapter', 'chap']) and chapter:
                shape.text_frame.text = chapter
                for paragraph in shape.text_frame.paragraphs:
                    paragraph.font.name = "Pretendard"
                    paragraph.font.size = Pt(18)
                modified_count += 1
            
            # Title text update
            elif any(keyword in text for keyword in ['title']) and title:
                shape.text_frame.text = title
                for paragraph in shape.text_frame.paragraphs:
                    paragraph.font.name = "Pretendard"
                    paragraph.font.size = Pt(24)
                modified_count += 1
                
            # Contents text update
            elif any(keyword in text for keyword in ['contents', 'content']) and contents:
                shape.text_frame.text = contents
                for paragraph in shape.text_frame.paragraphs:
                    paragraph.font.name = "Pretendard"
                    paragraph.font.size = Pt(14)
                modified_count += 1
        
        return f"""Slide {slide_number} updated successfully!

Modified elements: {modified_count}
Updated content:
   Chapter: {chapter if chapter else '(unchanged)'}
   Title: {title if title else '(unchanged)'}
   Contents: {contents[:50] + '...' if len(contents) > 50 else contents if contents else '(unchanged)'}"""
        
    except Exception as e:
        return f"Error updating slide {slide_number}: {str(e)}"

@mcp.tool()
def create_presentation_from_template(template_name: str, presentation_title: str = "New Presentation") -> str:
    """Create a new presentation from a template"""
    global current_presentation, current_filename
    
    if not template_registry:
        discover_templates()
    
    if template_name not in template_registry:
        available = ", ".join(template_registry.keys()) if template_registry else "None"
        return f"Template '{template_name}' not found.\nAvailable templates: {available}"
    
    try:
        current_presentation = load_template_presentation(template_name)
        
        if current_presentation is None:
            return f"Failed to load template '{template_name}'"
        
        current_filename = None
        
        # Update title slide if it exists
        if len(current_presentation.slides) > 0:
            title_slide = current_presentation.slides[0]
            if title_slide.shapes.title:
                title_slide.shapes.title.text = presentation_title
        
        template_info = template_registry[template_name]
        
        return f"""Presentation created from template!
Template: {template_name}
Source: {template_info['location']}
Title: {presentation_title}
Initial slides: {len(current_presentation.slides)}"""
        
    except Exception as e:
        return f"Error creating presentation from template: {str(e)}"

# ═══════════════════════════════════════════════════════════════════
# FILE MANAGEMENT TOOLS
# ═══════════════════════════════════════════════════════════════════

@mcp.tool()
def save_presentation(filename: Optional[str] = None, auto_save: bool = True) -> str:
    """Save presentation to local storage"""
    global current_presentation, current_filename
    
    if current_presentation is None:
        return "No presentation to save."
    
    try:
        if filename is None:
            if current_filename:
                filename = current_filename
            else:
                timestamp = datetime.datetime.now().strftime("%Y%m%d_%H%M%S")
                filename = f"presentation_{timestamp}.pptx"
        
        if not filename.endswith('.pptx'):
            filename += '.pptx'
        
        save_path = PRESENTATIONS_DIR / filename
        
        if save_path.exists():
            backup_name = f"{save_path.stem}_backup_{datetime.datetime.now().strftime('%H%M%S')}.pptx"
            backup_path = PRESENTATIONS_DIR / backup_name
            shutil.copy2(save_path, backup_path)
        
        current_presentation.save(str(save_path))
        current_filename = filename
        
        save_info = {
            "filename": filename,
            "path": str(save_path),
            "size": save_path.stat().st_size,
            "saved_at": datetime.datetime.now().isoformat(),
            "slide_count": len(current_presentation.slides)
        }
        
        if auto_save:
            meta_path = PRESENTATIONS_DIR / f"{save_path.stem}_meta.json"
            with open(meta_path, 'w', encoding='utf-8') as f:
                json.dump(save_info, f, indent=2, ensure_ascii=False)
        
        return f"""Presentation saved successfully!
Filename: {filename}
Path: {save_path}
Size: {save_info['size']:,} bytes
Slides: {save_info['slide_count']}
Saved at: {datetime.datetime.now().strftime('%Y-%m-%d %H:%M:%S')}"""
        
    except Exception as e:
        return f"Save error: {str(e)}"

@mcp.tool()
def get_presentation_info() -> str:
    """Get current presentation information"""
    global current_presentation, current_filename
    
    if current_presentation is None:
        return "No presentation open."
    
    try:
        slide_count = len(current_presentation.slides)
        filename = current_filename or "Not saved"
        
        slide_titles = []
        for i, slide in enumerate(current_presentation.slides, 1):
            title = "No title"
            if slide.shapes.title and slide.shapes.title.text:
                title = slide.shapes.title.text[:50]
            slide_titles.append(f"  {i}. {title}")
        
        return f"""Current presentation info:
Filename: {filename}
Slides: {slide_count}
Save directory: {PRESENTATIONS_DIR}
Status: Ready
Slide list:
{chr(10).join(slide_titles)}"""
        
    except Exception as e:
        return f"Info error: {str(e)}"

@mcp.tool()
def list_saved_presentations() -> str:
    """List all saved presentations"""
    try:
        pptx_files = list(PRESENTATIONS_DIR.glob("*.pptx"))
        
        if not pptx_files:
            return f"No saved presentations found.\nSave path: {PRESENTATIONS_DIR}"
        
        file_list = []
        for file_path in sorted(pptx_files, key=lambda x: x.stat().st_mtime, reverse=True):
            stat = file_path.stat()
            size_mb = stat.st_size / (1024 * 1024)
            modified = datetime.datetime.fromtimestamp(stat.st_mtime)
            
            file_list.append(f"{file_path.name}")
            file_list.append(f"   Size: {size_mb:.1f}MB")
            file_list.append(f"   Modified: {modified.strftime('%Y-%m-%d %H:%M:%S')}")
            file_list.append("")
        
        return f"""Saved presentations ({len(pptx_files)} files)
Save path: {PRESENTATIONS_DIR}

{chr(10).join(file_list)}"""
        
    except Exception as e:
        return f"Error listing files: {str(e)}"

@mcp.tool()
def create_presentation(title: str = "New Presentation") -> str:
    """Create a new PowerPoint presentation"""
    global current_presentation, current_filename
    
    try:
        current_presentation = Presentation()
        current_filename = None
        
        slide_layout = current_presentation.slide_layouts[0]
        slide = current_presentation.slides.add_slide(slide_layout)
        
        if slide.shapes.title:
            slide.shapes.title.text = title
            
        return f"New presentation created: '{title}'"
    except Exception as e:
        return f"Error creating presentation: {str(e)}"

@mcp.tool()
def add_slide(title: str = "", content: str = "", layout_index: int = 1) -> str:
    """Add new slide to presentation"""
    global current_presentation
    
    if current_presentation is None:
        return "No presentation open. Please create a presentation first."
    
    try:
        slide_layout = current_presentation.slide_layouts[layout_index]
        slide = current_presentation.slides.add_slide(slide_layout)
        
        if title and slide.shapes.title:
            slide.shapes.title.text = title
            
        if content and len(slide.placeholders) > 1:
            slide.placeholders[1].text = content
            
        slide_number = len(current_presentation.slides)
        return f"Added slide {slide_number}: '{title}'"
    except Exception as e:
        return f"Error adding slide: {str(e)}"

# ═══════════════════════════════════════════════════════════════════
# NOTION HELPER TOOLS
# ═══════════════════════════════════════════════════════════════════

@mcp.tool()
def validate_notion_structure(notion_url: str) -> str:
    """Notion 페이지 구조 검증"""
    validation = {
        'basic_settings': False,
        'slide_config': False,
        'style_guide': False,
        'has_tables': False,
        'has_diagrams': False
    }
    
    try:
        page_content = fetch_notion_page(notion_url)
        if page_content:
            content_text = page_content.get('text', '')
            
            validation['basic_settings'] = '기본 설정' in content_text
            validation['slide_config'] = '슬라이드 구성' in content_text  
            validation['style_guide'] = '스타일 가이드' in content_text
            validation['has_tables'] = '표 데이터' in content_text
            validation['has_diagrams'] = '다이어그램 요소' in content_text
            
        # 결과 포맷팅
        result_lines = ["📋 Notion 페이지 구조 검증 결과:"]
        
        for key, value in validation.items():
            status = "✅" if value else "❌"
            readable_key = {
                'basic_settings': '기본 설정',
                'slide_config': '슬라이드 구성',
                'style_guide': '스타일 가이드',
                'has_tables': '표 데이터',
                'has_diagrams': '다이어그램 요소'
            }.get(key, key)
            
            result_lines.append(f"  {status} {readable_key}")
        
        valid_count = sum(validation.values())
        result_lines.append(f"\n📊 유효한 섹션: {valid_count}/5")
        
        if valid_count >= 3:
            result_lines.append("🎉 PPT 자동 생성 가능!")
        else:
            result_lines.append("⚠️ 추가 설정이 필요합니다.")
        
        return "\n".join(result_lines)
    
    except Exception as e:
        return f"❌ 구조 검증 오류: {str(e)}"

@mcp.tool()
def get_slide_config_by_number(notion_url: str, slide_number: int) -> str:
    """특정 슬라이드 번호의 구성 정보 가져오기"""
    try:
        slide_configs = extract_slide_configurations(notion_url)
        for config in slide_configs:
            if int(config.get('슬라이드', 0)) == slide_number:
                return f"""슬라이드 {slide_number} 구성 정보:
                
Chapter: {config.get('Chapter', 'N/A')}
Title: {config.get('Title', 'N/A')}
Contents: {config.get('Contents', 'N/A')}
Layout Type: {config.get('Layout_Type', 'N/A')}
특별요구사항: {config.get('특별요구사항', 'N/A')}"""
        
        return f"❌ 슬라이드 {slide_number}의 구성 정보를 찾을 수 없습니다."
        
    except Exception as e:
        return f"❌ 슬라이드 구성 조회 오류: {str(e)}"

@mcp.tool()
def get_table_data_by_slide(notion_url: str, slide_number: int) -> str:
    """특정 슬라이드의 표 데이터 가져오기"""
    try:
        table_data = extract_table_data(notion_url)
        slide_tables = [item for item in table_data if item.get('Parent_Slide') == str(slide_number)]
        
        if not slide_tables:
            return f"❌ 슬라이드 {slide_number}에 표 데이터가 없습니다."
        
        # 표 데이터 정리해서 출력
        organized = organize_table_data(slide_tables)
        
        result_lines = [f"📊 슬라이드 {slide_number}의 표 데이터:"]
        
        for table_key, table_info in organized.items():
            result_lines.append(f"\n테이블 ID: {table_info['table_id']}")
            result_lines.append(f"크기: {table_info['max_row']}행 x {table_info['max_col']}열")
            
            # 셀 데이터 샘플 표시
            for cell in table_info['cells'][:5]:  # 처음 5개만
                result_lines.append(f"  ({cell['row']},{cell['col']}): {cell['value']} {'(헤더)' if cell['is_header'] else ''}")
            
            if len(table_info['cells']) > 5:
                result_lines.append(f"  ... 총 {len(table_info['cells'])}개 셀")
        
        return "\n".join(result_lines)
        
    except Exception as e:
        return f"❌ 표 데이터 조회 오류: {str(e)}"

# Initialize template discovery on startup
discover_templates()

if __name__ == "__main__":
    print(f"Enhanced PowerPoint MCP Server with Template Clone & Update + Notion Integration starting")
    print(f"Save directory: {PRESENTATIONS_DIR}")
    print(f"Temp directory: {TEMP_DIR}")
    print(f"Template registry: {len(template_registry)} templates found")
    mcp.run()